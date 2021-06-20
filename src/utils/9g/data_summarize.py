import csv
import datetime
import os
import re
import sys
import time
from collections import OrderedDict
from glob import glob

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import win32com.client
from openpyxl import load_workbook
from openpyxl.chart import LineChart, Reference
from openpyxl.chart.layout import Layout, ManualLayout
from openpyxl.chart.shapes import GraphicalProperties
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from variables import CELL_WIDTH_BASE
from variables import CELL_WIDTH_BASE_PIN
from variables import CELL_WIDTH_BASE_RATE
from variables import CELL_WIDTH_BASE_VI
from variables import DATA_GROUP
from variables import DATA_INDEX
from variables import FREQ_YTICKS

# from variables import CROSSTALK_FILE_NAME
# from variables import HISTOGRAM_FILE_NAME
from variables import OVERVIEW_FILE_NAME

# from variables import OVERVIEW_FILE_NAME
from variables import RENAME_CONDITIONS

# from openpyxl.drawing.line import LineProperties

picture_counter = 0

now = datetime.datetime.now()
date_now = now.strftime("%Y%m%d%H%M")

"""
TODO
# TODO fix this code, change to easy to read file name
title = file.replace("\\", "xyz").replace(".png", "")
title = re.sub(".*xyz", "", title)

"""

"""docstring sample. function simple explanation is written in this line.

    if function detail explanation is necessary, write to here.

    Args:
        param1 (int): The first parameter
        param2 (str): The second parameter

    Returns:
        bool: The return value. True for success, False otherwise

    https://google.github.io/styleguide/pyguide.html
"""


def check_pin_kind(pin_name):
    """check pin kind and return pin kind and order for sort

    Args:
        pin_name(str): pin name like P1857A1

    Returns:
        pin_kind_for_pptx, pin_order
        pin_kind_for_pptx: "IO", "WCK", "CK", "CA", "CS",
        pin_order: "IO"->1, "WCK"->2, "CK"->3, "CA"->4, "CS"->5

    """
    match_pin_name = re.match(r"P(\d*).*", pin_name)
    pin_num = int(match_pin_name.group(1))

    if pin_num < 1857:
        pin_kind = "IO"
        pin_order = 1
    elif 1857 <= pin_num <= 1888:
        pin_kind = "WCK"
        pin_order = 2
    elif 1889 <= pin_num <= 1890:
        pin_kind = "CK"
        pin_order = 3
    elif 1921 <= pin_num <= 1933:
        pin_kind = "CA"
        pin_order = 4
    elif 1953 <= pin_num <= 1959:
        pin_kind = "CS"
        pin_order = 5
    else:
        print("Pin_kind Error")
        sys.exit()

    return pin_kind, pin_order


def mul3(x):
    return x * 1e3


def mul12(x):
    return x * 1e12


def mul_m9(x):
    return x * 1e-9


def get_nearest_value(lst, num):
    """return nearest value of num from list

    Args:
        lst (list): list of num
        num (int): num of value

    Returns:
        nearest value of num from list

    """
    idx = np.abs(np.asarray(lst) - num).argmin()
    return lst[idx]


def wf_txt_data_to_csv(file):
    """create csv file from osc output text file

    Args:
        file (str): file name

    Returns:
        None

    """
    print(file)
    with open(file.replace(".txt", ".csv"), "w", encoding="utf-8") as fw:
        flg = 0
        with open(file, encoding="utf-8") as fr:
            for line in fr.read().splitlines():
                match_data = re.match(r"Data", line)

                if line != "":
                    if flg:
                        fw.write(line)
                        fw.write("\n")

                    if match_data:
                        flg = 1


class WaveData:
    def __init__(
        self,
        file_name,
        folder_path,
        active_presentation,
        pptx_lib,
        index="Pin_Rate",
        group_by=None,
        header=None,
    ):
        self.data_df = pd.DataFrame()
        self.data_vns = []
        self.data_vix = []
        self.data_overshoot = []
        self.file_name = file_name
        self.folder_path = folder_path
        self.group_by = group_by
        self.header = header
        self.index = index
        self.input_file = self.folder_path + self.file_name
        self.pptx_lib = pptx_lib
        self.active_presentation = active_presentation
        self.chart = None
        self.fig = None
        self.ax = None
        self.slide = None

        if self.pptx_lib == "win32com":
            self.slide_width = self.active_presentation.PageSetup.SlideWidth
            self.slide_height = self.active_presentation.PageSetup.SlideHeight
            self.slide_count = self.active_presentation.Slides.Count

        elif self.pptx_lib == "python-pptx":
            self.slide_width = self.active_presentation.slide_width
            self.slide_height = self.active_presentation.slide_height
            self.slide_count = 0

        self.make_df_and_xlsx()

    def update_slide_count(self):
        if self.pptx_lib == "win32com":
            self.slide_count = self.active_presentation.Slides.Count

        elif self.pptx_lib == "python-pptx":
            self.slide_count = 0

    def make_df_and_xlsx(self):
        """Make pandas dataframe and xlsx data from csv file

        Args:

        Returns:
            None

        """
        with open(self.input_file, mode="r", encoding="utf-8-sig") as csv_file:
            reader = csv.reader(csv_file)
            data = []

            for rows in reader:
                match = re.match(r"(P(\d*).*?)_.*?(Vih.*?_.*?_.*?)_(.*?)_.*", rows[0])
                if match:
                    rows.insert(0, "Condition")
                    rows.insert(2, "Pin")
                    rows.insert(3, match.group(1))
                    rows.insert(4, "Pin_kind")

                    pin_kind, pin_order = check_pin_kind(match.group(1))

                    rows.insert(5, pin_kind)
                    rows.insert(6, "Vi")
                    rows.insert(7, match.group(3).replace("00V", "V"))
                    rows.insert(8, "Rate")
                    rows.insert(
                        9, match.group(4).replace("Rate0r", "").replace("ns", "ps")
                    )
                    rows.insert(10, "Order")
                    rows.insert(11, str(pin_order))
                    rows.insert(12, "Condition_all")
                    rows.insert(13, match.group(0))

                dic = OrderedDict()
                for i in range(0, len(rows) - 1, 2):
                    try:
                        dic[rows[i].replace(" ", "").capitalize()] = float(
                            rows[i + 1].replace(" ", "")
                        )
                    except ValueError:
                        dic[rows[i].replace(" ", "").capitalize()] = rows[
                            i + 1
                        ].replace(" ", "")
                    except AttributeError:
                        dic[rows[i].replace(" ", "").capitalize()] = rows[i + 1]

                data.append(dic)

            self.data_df = pd.DataFrame(data)

            self.data_df.insert(
                1, "Pin_Rate", self.data_df["Pin"] + "_" + self.data_df["Rate"]
            )
            self.data_df.insert(
                2, "Pin_Vi", self.data_df["Pin"] + "_" + self.data_df["Vi"]
            )
            self.data_df.insert(
                3, "Pin_kind_Vi", self.data_df["Pin_kind"] + "_" + self.data_df["Vi"]
            )

            # adjust unit of dataframe
            self.adjust_unit()

            if self.header:
                self.data_df = self.data_df.set_axis(self.header, axis="columns")

            self.data_df = self.data_df.set_index(self.index)

            with pd.ExcelWriter(self.input_file.replace("csv", "xlsx")) as writer:
                self.data_df.to_excel(writer, sheet_name="summary")

                if self.group_by:
                    for name, group in self.data_df.groupby(self.group_by):
                        group.to_excel(writer, sheet_name=name)

    def make_excel_graph(
        self,
        file_name,
        chart_yaxis_scaling,
        chart_height=9,
        chart_width=16,
        chart_position="C2",
        chart_yaxis_title=None,
    ):
        """make specified excel graph using xlsx data

        Args:
            file_name (str): input excel file name
            chart_yaxis_scaling (list): values for yaxis scale.
            chart_height (float): chart height.
            chart_width (float): chart width.
            chart_position (str): chart position at excel.
            chart_yaxis_title: chart yaxis title

        Returns:
            None

        """

        wb = load_workbook(file_name)
        for i in range(len(wb.worksheets)):
            ws = wb.worksheets[i]

            values = Reference(
                ws, min_row=1, min_col=2, max_row=ws.max_row, max_col=ws.max_column,
            )
            categories = Reference(
                ws, min_row=2, min_col=1, max_row=ws.max_row, max_col=1
            )

            self.setup_excel_chart(
                values=values,
                categories=categories,
                chart_height=chart_height,
                chart_width=chart_width,
                chart_yaxis_title=chart_yaxis_title,
                chart_yaxis_scaling_min=chart_yaxis_scaling[0],
                chart_yaxis_scaling_max=chart_yaxis_scaling[1],
                chart_yaxis_major_unit=chart_yaxis_scaling[2],
            )

            ws.add_chart(self.chart, chart_position)
        wb.save(file_name)

    def make_excel_graphs(
        self,
        chart_height=9,
        chart_width=16,
        chart_yaxis_titles=None,
        chart_yaxis_scaling_mines=None,
        chart_yaxis_scaling_maxes=None,
        chart_yaxis_major_unit=None,
        data_start_column=0,
    ):
        """make excel graphs using xlsx data"""

        wb = load_workbook(self.input_file.replace("csv", "xlsx"))
        for i in range(len(wb.worksheets)):
            ws = wb.worksheets[i]

            for j in range(ws.max_column + data_start_column * -1 + 1):
                values = Reference(
                    ws,
                    min_row=1,
                    min_col=data_start_column + j,
                    max_row=ws.max_row,
                    max_col=data_start_column + j,
                )
                categories = Reference(
                    ws, min_row=2, min_col=1, max_row=ws.max_row, max_col=1
                )
                self.setup_excel_chart(
                    values=values,
                    categories=categories,
                    chart_height=chart_height,
                    chart_width=chart_width,
                    chart_yaxis_title=chart_yaxis_titles[j],
                    chart_yaxis_scaling_min=chart_yaxis_scaling_mines[j],
                    chart_yaxis_scaling_max=chart_yaxis_scaling_maxes[j],
                    chart_yaxis_major_unit=chart_yaxis_major_unit[j],
                )

                ws.add_chart(self.chart, "B" + str(5 + 20 * j))
        wb.save(self.input_file.replace("csv", "xlsx"))

    def setup_excel_chart(
        self,
        values,
        categories,
        chart_height,
        chart_width,
        chart_yaxis_title,
        chart_yaxis_scaling_min,
        chart_yaxis_scaling_max,
        chart_yaxis_major_unit,
    ):
        """excel chart setup

        Args:
            values (list): values for excel graph.
            categories (list): index values for excel graph.
            chart_height (float): chart height
            chart_width (float): chart width
            chart_yaxis_title (str): chart yaxis title
            chart_yaxis_scaling_min (float): chart yaxis min value
            chart_yaxis_scaling_max (float): chart yaxis max value
            chart_yaxis_major_unit (float): chart yaxis major unit

        Returns:
            None

        """

        self.chart = LineChart()
        self.chart.add_data(values, titles_from_data=True)
        self.chart.title = ""
        self.chart.set_categories(categories)
        self.chart.height = chart_height
        self.chart.width = chart_width
        self.chart.y_axis.title = chart_yaxis_title
        self.chart.y_axis.scaling.min = chart_yaxis_scaling_min
        self.chart.y_axis.scaling.max = chart_yaxis_scaling_max
        self.chart.y_axis.majorUnit = chart_yaxis_major_unit
        self.chart.y_axis.numFmt = "0.0"
        self.chart.x_axis.tickLblPos = "low"
        self.chart.layout = Layout(
            ManualLayout(x=0.1, y=0.1, h=0.8, w=0.8, xMode="edge", yMode="edge")
        )
        # self.chart.legend.legendPos = "tr"
        self.chart.legend.layout = Layout(
            manualLayout=ManualLayout(yMode="edge", xMode="edge", x=0.75, y=0.13)
        )
        self.chart.legend.spPr = GraphicalProperties()
        self.chart.legend.spPr.solidFill = "FFFFFF"
        self.chart.legend.spPr.ln.solidFill = "E0E0E0"
        self.chart.legend.spPr.ln.w = 1 * 12700

        # chart.y_axis.majorGridlines = None

        for i in range(len(self.chart.series)):
            series = self.chart.series[i]
            series.marker.symbol = "circle"
            # series.marker.size = 10
            series.graphicalProperties.line.noFill = True

    def make_graph(
        self,
        df_columns_list,
        y_ticks,
        figure_size=(10, 5.5),
        file_name="default",
        font_size=14,
        digit_format="%.1f",
        legend_loc="upper right",
        rotation=45,
        styles=None,
        ax_h_lines=None,
        spec=False,
        legends=None,
        pin_kind=None,
        x_label=None,
        y_label=None,
        y_ticks_per_condition=None,
        ax_h_lines_per_condition=None,
        additional_information=False,
        info=None,
    ):
        """make specified graph from dataframe using matplotlib

        Args:
            df_columns_list (list): dataframe columns list to make graph
            y_ticks (list): y_ticks
            figure_size (tuple): figure size
            file_name (str): filename
            font_size (int): font size
            digit_format (str): axis digit_format setting
            legend_loc (str): legend location
            rotation (int): rotation
            styles (list): marker style
            ax_h_lines (list): yaxis reference line
            spec (bool): set true if spec condition
            legends (dict): legend
            pin_kind (str): pin kind
            x_label (str): x_label
            y_label (str): y_label
            y_ticks_per_condition (dict):
            ax_h_lines_per_condition:
            additional_information:
            info:

        Returns:
            None

        """
        global picture_counter

        # for excel graph
        os.makedirs(self.folder_path + "excel_graph_data", exist_ok=True)

        # if needs to separate result per pin kind
        if pin_kind:
            df = self.data_df[self.data_df["Pin_kind"] == pin_kind].copy()
        else:
            df = self.data_df.copy()

        y_ticks_tmp = y_ticks  # save original y_ticks setting
        if self.group_by:
            for name, group in df.groupby(self.group_by, sort=False):
                print(f"\ncondition_name = {name}")
                print(group)
                df_plot = group[df_columns_list]
                print(df_plot)

                # skip if dataframe has missing value
                if df_plot.isnull().values.sum() != 0:
                    df_plot = group[df_columns_list].dropna(how="all")
                    print(df_plot)

                    if df_plot.empty:
                        continue

                index_count = len(df_plot.index)

                if index_count == 2:
                    x_margin = 0.5

                else:
                    x_margin = 0.1

                self.setup_fig_and_ax(
                    figure_size=figure_size,
                    bottom=0.3,
                    x_margin=x_margin,
                    digit_format=digit_format,
                )

                # set number of label
                self.ax.set_xticks([i for i in range(group.shape[0])])

                y_ticks = y_ticks_tmp
                if y_ticks_per_condition and (name in y_ticks_per_condition):
                    y_ticks = y_ticks_per_condition[name]

                print(df_plot)

                if legends is not None:
                    df_plot = df_plot.rename(columns=legends, inplace=False)

                if styles is None:
                    styles = ["o", "o", "o", "o"]

                df_plot.plot(
                    ax=self.ax,
                    ylim=y_ticks[:2],
                    style=styles,
                    legend=True,
                    fontsize=font_size,
                )

                if additional_information:
                    self.add_ax_text(
                        x=0.99,
                        y=0.05,
                        s=info,
                        transform=self.ax.transAxes,
                        horizontal_alignment="right",
                    )

                self.adjust_graph_params(
                    group_name=str(name),
                    rotation=rotation,
                    x_label=x_label,
                    y_label=y_label,
                    y_ticks=y_ticks,
                    font_size=font_size,
                    ax_h_lines=ax_h_lines,
                    legend_loc=legend_loc,
                    spec=spec,
                    grid=True,
                    ax_h_lines_per_condition=ax_h_lines_per_condition,
                )

                picture_number = f"{picture_counter:03}_"

                # for excel graph
                excel_file_name = (
                    self.folder_path
                    + "/excel_graph_data/"
                    + picture_number
                    + self.file_name.replace(".csv", "")
                    + "_"
                    + name
                    + "_"
                    + "_".join(df_columns_list)
                    + ".xlsx"
                )

                df_to_excel = df_plot.set_axis(legends, axis=1)
                df_to_excel.to_excel(excel_file_name)
                self.make_excel_graph(
                    file_name=excel_file_name,
                    chart_yaxis_scaling=y_ticks,
                    chart_yaxis_title=y_label,
                )

                file_path = (
                    self.folder_path + picture_number + name + "_" + file_name + ".png"
                )
                plt.savefig(file_path)
                plt.close("all")

                print(name)

                for key, value in RENAME_CONDITIONS.items():
                    name = name.replace(key, value)

                file_name = file_name.replace("_", " ")
                print(file_name)
                # self.add_slide_to_pptx(
                #     title=picture_number + name + " " + file_name,
                #     layout=11,
                # )
                self.add_slide_to_pptx(
                    title=file_name + " " + name, layout=11,
                )

                self.add_picture_to_pptx(file_path=file_path)

        else:
            self.setup_fig_and_ax(figure_size, bottom=0.3)

            # set number of label
            self.ax.set_xticks([i for i in range(self.data_df.shape[0])])

            self.data_df[df_columns_list].plot(
                ax=self.ax,
                ylim=y_ticks[:2],
                style=styles,
                legend=True,
                fontsize=font_size,
            )

            self.adjust_graph_params(
                rotation=rotation,
                x_label=x_label,
                y_label=y_label,
                font_size=font_size,
                y_ticks=y_ticks,
                ax_h_lines=ax_h_lines,
            )

            picture_number = f"{picture_counter:03}_"
            file_path = self.folder_path + picture_number + file_name + ".png"
            plt.savefig(file_path)
            plt.close("all")
            self.add_slide_to_pptx(
                title=picture_number + file_name, layout=11,
            )

            self.add_picture_to_pptx(file_path=file_path)

    def make_overshoot_graph(
        self,
        file,
        y_ticks,
        figure_size=(10, 5.5),
        item_name="Overshoot",
        additional_information=False,
        info=None,
    ):
        """make overshoot graph

        Args:
            file (str): waveform text data to make graph
            figure_size (tuple): figure size
            item_name (str): item name. Overshoot or Undershoot
            additional_information (bool):  additional information flag
            info (str): additional information

        Returns:
            None
            :param info:
            :param additional_information:
            :param item_name:
            :param figure_size:
            :param file:
            :param y_ticks:

        """
        self.setup_fig_and_ax(
            figure_size=figure_size, x_margin=0.01, digit_format="%.3f"
        )

        match_pin_file = re.match(r".*((P.*?)_.*(Vih.*)_(Rate0r.*ns).*).txt", file)
        condition_all = match_pin_file.group(1)
        print(self.data_df)
        print(
            self.data_df[
                (self.data_df["Condition_all"] == condition_all)
                & self.data_df[item_name].notna()
            ]
        )
        df_tmp = self.data_df[
            (self.data_df["Condition_all"] == condition_all)
            & self.data_df[item_name].notna()
        ]
        if (
            item_name == "Overshoot"
            and df_tmp["Top"].size == 1
            and df_tmp["Vmaximum"].size == 1
        ):
            reference_level = df_tmp["Top"][0] * 1e-3
            vmaximum = df_tmp["Vmaximum"][0] * 1e-3
        elif (
            item_name == "Undershoot"
            and df_tmp["Base"].size == 1
            and df_tmp["Vminimum"].size == 1
        ):
            reference_level = df_tmp["Base"][0] * 1e-3
            vmaximum = df_tmp["Vminimum"][0] * 1e-3
        else:
            print(
                "Over 2 results/No result in result file for overshoot/undershoot graph"
            )
            sys.exit()

        print(reference_level)
        print(vmaximum)

        pin_name = match_pin_file.group(2)
        test_rate = match_pin_file.group(4)
        vi = match_pin_file.group(3).replace("00V", "V")

        wf_txt_data_to_csv(file)

        df = pd.read_csv(file.replace(".txt", ".csv"), header=None)
        df = df.set_axis(["t", pin_name], axis=1)
        df = df.set_index("t")

        x = np.array(df.index.tolist())
        y = np.array(df[pin_name].tolist())

        graph_x_middle = int(x.size / 2)

        sum_of_voltage = 0
        x_start = 0
        x_start_flag = 0
        x_end = 0
        area_label_y_position = 0
        overshoot_label_y_position = 0
        if item_name == "Overshoot":
            y_label_position_offset = 0.075
            area_label_y_position = vmaximum + y_label_position_offset
            overshoot_label_y_position = vmaximum + y_label_position_offset * 2
            self.ax.fill_between(
                x=x,
                y1=y,
                y2=reference_level,
                where=(y > reference_level) & (x < x[graph_x_middle]),
                color="C0",
                alpha=0.2,
            )
            for i in range(int(x.size / 2)):
                if y[i] > reference_level:
                    if x_start_flag == 0:
                        x_start = x[i]
                        x_start_flag += 1

                    sum_of_voltage += y[i] - reference_level
                    x_end = x[i]

        elif item_name == "Undershoot":
            y_label_position_offset = -0.075
            area_label_y_position = vmaximum + y_label_position_offset * 2
            overshoot_label_y_position = vmaximum + y_label_position_offset
            self.ax.fill_between(
                x=x,
                y1=y,
                y2=reference_level,
                where=(y < reference_level) & (x < x[graph_x_middle]),
                color="C0",
                alpha=0.2,
            )
            for i in range(int(x.size / 2)):
                if y[i] < reference_level:
                    if x_start_flag == 0:
                        x_start = x[i]
                        x_start_flag += 1

                    sum_of_voltage += y[i] - reference_level
                    x_end = x[i]

        ratio_of_area_per_1ns = (x[2047] - x[0] + (x[1] - x[0])) * 1e9 / x.size
        area = abs(sum_of_voltage * ratio_of_area_per_1ns)
        overshoot = abs(vmaximum - reference_level) * 1e3
        s = "area"
        self.add_ax_text(
            x=(x_start + x_end) / 2,
            y=area_label_y_position,
            s=f"{s:10} = {area:.6f}[V-ns]",
            transform=self.ax.transData,
            z_order=11,
            horizontal_alignment="left",
        )
        self.add_ax_text(
            x=(x_start + x_end) / 2,
            y=overshoot_label_y_position,
            s=f"{item_name:10} = {overshoot:.6f}[mV]",
            transform=self.ax.transData,
            z_order=12,
            horizontal_alignment="left",
        )
        vns_data = {
            "Condition": condition_all,
            "v-ns-" + item_name: area,
            item_name + "_mV": overshoot,
        }
        self.data_vns.append(vns_data)

        df.plot(ax=self.ax, ylim=y_ticks[:2])

        if additional_information:
            self.add_ax_text(
                x=0.99,
                y=0.05,
                s=info,
                transform=self.ax.transAxes,
                horizontal_alignment="right",
            )

        # make data for table output
        self.data_overshoot.append(
            {"Vi": vi, "Pin": pin_name, "rate": test_rate, "overshoot(v-ns)": 0}
        )

        self.adjust_graph_params(
            rotation=0,
            x_label="",
            y_label="mV",
            # font_size=font_size,
            ax_h_lines=[reference_level, vmaximum],
            y_ticks=y_ticks,
        )
        picture_number = f"{picture_counter:03}_"
        pin_kind = check_pin_kind(pin_name)
        file_path = (
            self.folder_path
            + picture_number
            + pin_kind[0]
            + "_"
            + vi
            + "_"
            + item_name
            + ".png"
        )
        plt.savefig(file_path)
        plt.close("all")
        title = item_name + "_ " + pin_kind[0] + "_" + vi + "_" + test_rate
        for key, value in RENAME_CONDITIONS.items():
            title = title.replace(key, value)

        self.add_slide_to_pptx(
            title=title, layout=11,
        )
        self.add_picture_to_pptx(file_path=file_path)

    def make_vix_graph(
        self,
        item_name,
        negative_pin_file,
        positive_pin_file,
        description=False,
        font_size=14,
        figure_size=(10, 5.5),
        reference_level=0,
        rotation=0,
        x_label=None,
        y_label=None,
    ):
        """make vix graph from positive/negative wave data file using matplotlib

        Args:
            item_name (str): item name (Vix)
            negative_pin_file (str): csv negative pin file name
            positive_pin_file (str): csv positive pin file name
            description (bool): if put vix min/max ft description
            font_size (int): font_size
            figure_size (tuple): figure size
            reference_level (float): reference level
            rotation (int): x_label rotation value
            x_label (str): x_label
            y_label (str): y_label

        Returns:
            None

        """
        global picture_counter

        self.setup_fig_and_ax(figure_size=figure_size, x_margin=0.01)

        match_positive_pin = re.match(
            r".*(P.*?)_.*(Vih.*)_Rate0r(.*ns).*", positive_pin_file
        )
        match_negative_pin = re.match(
            r".*(P.*?)_.*(Vih.*)_Rate0r(.*ns).*", negative_pin_file
        )

        positive_pin_name = match_positive_pin.group(1)
        negative_pin_name = match_negative_pin.group(1)
        test_rate = match_positive_pin.group(3)
        vi = match_positive_pin.group(2).replace("00V", "V")

        wf_txt_data_to_csv(positive_pin_file)
        wf_txt_data_to_csv(negative_pin_file)

        df_positive = pd.read_csv(
            positive_pin_file.replace(".txt", ".csv"), header=None
        )
        df_negative = pd.read_csv(
            negative_pin_file.replace(".txt", ".csv"), header=None
        )

        df_positive = df_positive.set_axis(["t", "wck_t"], axis=1)
        df_negative = df_negative.set_axis(["t", "wck_c"], axis=1)

        df_positive = df_positive.set_index("t")
        df_negative = df_negative.set_index("t")

        df_positive_negative = pd.concat([df_positive, df_negative], axis=1)

        # make diff column
        df_positive_negative["f(t)"] = (
            df_positive_negative["wck_t"] - df_positive_negative["wck_c"]
        )

        # get 1 cycle waveform
        df_positive_negative = df_positive_negative.iloc[
            int(len(df_positive_negative) * 0.23) : int(
                len(df_positive_negative) * 0.73
            ),
            :,
        ]

        # make dataframe df_vix. df_vix has 2 data which wck_t - wck_c is close to 0
        # close to 0 or 0 means cross point
        df_tmp = df_positive_negative.copy()
        df_vix = pd.DataFrame()
        cross_point_count = 2
        for _ in range(cross_point_count):
            val = get_nearest_value(df_tmp["f(t)"].values.tolist(), 0)
            min_row1 = df_tmp[df_tmp["f(t)"] == val]
            df_vix = pd.concat([df_vix, min_row1])
            df_tmp = df_tmp.drop(min_row1.index)

        # get average in case there is no cross point in data
        df_vix["(wck_t+wck_c)/2"] = (df_vix["wck_t"] + df_vix["wck_c"]) / 2

        df_vix = df_vix["(wck_t+wck_c)/2"]
        df_vix = df_vix.reset_index()
        df_vix_list = df_vix.values.tolist()

        # add x, y coordinates of differential input cross point voltage to graph
        x_position_offset = 0
        y_position_offset = 0
        vix_wck_rf = 0
        vix_wck_fr = 0
        for df_vix_p in df_vix_list:
            x_position = df_vix_p[0]
            y_position = df_vix_p[1]

            if (
                x_position
                < df_positive_negative.index[int(len(df_positive_negative) / 2)]
            ):
                label = "Vix_WCK_FR"
                vix_wck_fr = y_position - reference_level
            else:
                label = "Vix_WCK_RF"
                vix_wck_rf = y_position - reference_level

            x_position_offset = 0.0005e-8
            y_position_offset = 0.05

            self.add_ax_text(
                x=x_position + x_position_offset,
                y=(y_position + reference_level) / 2,
                s=f"{label}={y_position - reference_level:.2f}mV",
                transform=self.ax.transData,
                z_order=11,
                horizontal_alignment="left",
            )
            self.add_ax_annotate(
                text="",
                xy=[x_position, y_position],
                xy_text=[x_position, reference_level],
                arrow_style="<->",
                z_order=10,
            )

        # for Min(f(t)), Max(f(t))
        max_index = df_positive_negative["f(t)"].idxmax()
        min_index = df_positive_negative["f(t)"].idxmin()
        max_index_values = df_positive_negative.loc[max_index]
        min_index_values = df_positive_negative.loc[min_index]

        # Max(f(t))
        self.add_ax_annotate(
            text="",
            xy=[max_index_values.name, max_index_values["wck_t"]],
            xy_text=[max_index_values.name, max_index_values["wck_c"]],
            arrow_style="->",
        )

        max_ft = max_index_values["wck_t"] - max_index_values["wck_c"]
        self.add_ax_text(
            x=max_index_values.name + x_position_offset,  # includes offset
            y=(max_index_values["wck_t"] + max_index_values["wck_c"]) / 2
            + y_position_offset,
            s=f"Max(f(t))={max_ft:.2f}mV",
            transform=self.ax.transData,
            horizontal_alignment="left",
        )

        # Min(f(t))
        self.add_ax_annotate(
            text="",
            xy=[min_index_values.name, min_index_values["wck_t"]],
            xy_text=[min_index_values.name, min_index_values["wck_c"]],
            arrow_style="->",
        )

        min_ft = min_index_values["wck_t"] - min_index_values["wck_c"]
        self.add_ax_text(
            x=min_index_values.name + x_position_offset,  # includes offset
            y=(min_index_values["wck_t"] + min_index_values["wck_c"]) / 2
            + y_position_offset,
            s=f"Min(f(t))={min_ft:.2f}mV",
            transform=self.ax.transData,
            z_order=11,
            horizontal_alignment="left",
        )

        # Vix_WCK_Ratio Calculation result
        x_position_vix_ratio_result = 0.25
        vix_wck_ratio_fr_min_t = (vix_wck_rf / abs(max_ft)) * 100
        vix_wck_ratio_rf_max_t = (vix_wck_fr / abs(min_ft)) * 100
        self.add_ax_text(
            x=x_position_vix_ratio_result,
            y=-0.2,
            s=f"Vix_WCK_Ratio = Vix_WCK_FR/|Min(f(t))| = {vix_wck_fr:5.2f}/|{min_ft:5.2f}| = {vix_wck_ratio_rf_max_t:4.1f}%",
            transform=self.ax.transAxes,
            horizontal_alignment="left",
        )
        self.add_ax_text(
            x=x_position_vix_ratio_result,
            y=-0.25,
            s=f"Vix_WCK_Ratio = Vix_WCK_Rf/ Max(f(t))  = {vix_wck_rf:.2f}/ {max_ft:5.2f}  = {vix_wck_ratio_fr_min_t:4.1f}%",
            transform=self.ax.transAxes,
            horizontal_alignment="left",
        )

        # make data for table output
        self.data_vix.append(
            {
                "Vi": vi,
                "Positive Pin": positive_pin_name,
                "Negative Pin": negative_pin_name,
                "rate": test_rate,
                "Vix_WCK_FR/|Min(f(t))| (%)": vix_wck_ratio_fr_min_t,
                "Vix_WCK_Rf/Max(f(t)) (%)": vix_wck_ratio_rf_max_t,
            }
        )

        # reference level line
        self.ax.hlines(
            y=reference_level,
            xmin=df_positive_negative.index[0],
            xmax=df_positive_negative.index[len(df_positive_negative) - 1],
            color="black",
            linestyle="dashed",
            zorder=10,
        )

        df_positive_negative = df_positive_negative.drop("f(t)", axis=1)
        df_positive_negative = df_positive_negative.rename(
            columns={"wck_t": positive_pin_name, "wck_c": negative_pin_name},
            inplace=False,
        )

        df_positive_negative.plot(ax=self.ax)

        self.adjust_graph_params(
            rotation=rotation,
            x_label=x_label,
            y_label=y_label,
            font_size=font_size,
            y_ticks=[],
            ax_h_lines=[],
        )
        picture_number = f"{picture_counter:03}_"
        pin_kind = check_pin_kind(positive_pin_name)
        file_path = (
            self.folder_path
            + picture_number
            + pin_kind[0]
            + "_"
            + vi
            + "_"
            + item_name
            + ".png"
        )
        plt.savefig(file_path)
        plt.close("all")
        self.add_slide_to_pptx(
            title=picture_number + pin_kind[0] + "_" + vi + "_" + item_name, layout=11,
        )
        self.add_picture_to_pptx(file_path=file_path)

        # insert Min(f(t)), Max(f(t)), Vix example pic from spec sheet
        if description:
            self.add_slide_to_pptx(
                title="Vix", layout=11,
            )
            self.add_picture_to_pptx(
                file_path=os.getcwd() + "/pictures/Vix.png",
                resize=True,
                count_picture=False,
            )
            self.add_slide_to_pptx(
                title="Min(f(t)), Max(f(t))", layout=11,
            )
            self.add_picture_to_pptx(
                file_path=os.getcwd() + "/pictures/ft.png",
                resize=True,
                count_picture=False,
            )

    def add_ax_annotate(self, text, xy, xy_text, arrow_style="<->", z_order=10):
        self.ax.annotate(
            text=text,
            xy=xy,
            xytext=xy_text,
            arrowprops=dict(arrowstyle=arrow_style),
            zorder=z_order,
            size=5,
        )

    def add_ax_text(
        self, x, y, s, transform, z_order=10, horizontal_alignment="center"
    ):
        self.ax.text(
            x=x,
            y=y,
            s=s,
            transform=transform,
            backgroundcolor="white",
            fontfamily="monospace",
            zorder=z_order,
            horizontalalignment=horizontal_alignment,
        )

    def setup_fig_and_ax(
        self,
        figure_size=(16, 9),
        top=0.95,
        left=0.15,
        bottom=0.2,
        right=0.85,
        x_margin=0.1,
        digit_format="%.1f",
    ):
        """set up matplotlib fix and ax object

        Args:
            figure_size (tuple): fig size
            bottom (float): bottom margin
            top (float): top margin
            left (float):  left margin
            right (float):  1 - right margin
            x_margin (float): x_margin
            digit_format (str): yaxis digit_format setting

        Returns:
            None

        """
        self.fig = plt.figure(figsize=figure_size)  # create figure object
        self.ax = self.fig.add_subplot(1, 1, 1, xmargin=x_margin)  # create axes object
        self.ax.yaxis.set_major_formatter(plt.FormatStrFormatter(digit_format))
        self.fig.subplots_adjust(top=top, left=left, bottom=bottom, right=right)

    def adjust_graph_params(
        self,
        # legends,
        y_ticks,
        font_size=14,
        legend_loc="upper right",
        rotation=0,
        group_name="",
        ax_h_lines=None,
        x_label=None,
        y_label=None,
        spec=False,
        grid=False,
        ax_h_lines_per_condition=None,
    ):
        """adjust graph parameters

        Args:
            y_ticks (list): y_ticks min, max, resolution
            font_size (int): font size
            legend_loc (str): legend location
            rotation (int): rotation
            group_name (str): group name
            ax_h_lines (list): ax_h_lines
            x_label (str): x_label
            y_label (str): y_label
            spec (bool): set true if spec condition
            grid (bool): grid
            ax_h_lines_per_condition (dict): ax h line per condition

        Returns:
            None

        """
        plt.xticks(rotation=rotation)
        self.ax.set_ylabel(y_label, fontsize=font_size)
        self.ax.set_xlabel(x_label, fontsize=font_size)
        self.ax.legend(
            # labels=legends,
            fontsize=font_size,
            loc=legend_loc,
            frameon=True,
            framealpha=1.0,
        )

        # set grid
        if grid:
            self.ax.grid(axis="y", linestyle="-", color="black", linewidth=1, alpha=0.2)

        if y_ticks:
            self.ax.set_yticks(
                np.arange(y_ticks[0], y_ticks[1] + y_ticks[2], y_ticks[2])
            )

        # add limit line in case AT condition Vih/Vil=1V/0V
        match_at_condition = re.match(r".*Vih1r0V_Vil0r0V", group_name)
        if spec and match_at_condition:
            line_style = "-"
            alpha = 0.8
        else:
            line_style = "--"
            alpha = 0.5

        if (
            ax_h_lines_per_condition is not None
            and group_name in ax_h_lines_per_condition
        ):
            ax_h_lines = ax_h_lines_per_condition[group_name]

        for ax_h_line in ax_h_lines:
            self.ax.axhline(
                y=ax_h_line,
                linestyle=line_style,
                alpha=alpha,
                color="gray",
                linewidth=1,
            )

    def add_vix_table_to_pptx(self, title, items, cell_width, cell_height=20):
        """add vix table to pptx

        Args:
            title (str): slide title
            items (list): items for table
            cell_width (list): cell width
            cell_height (int): cell height

        Returns:
            None

        """
        vix_data_list_to_table_df = pd.DataFrame(self.data_vix)
        print(vix_data_list_to_table_df)
        self.add_slide_to_pptx(title=title, layout=4)
        self.add_table(
            df=vix_data_list_to_table_df,
            items=items,
            cell_width=cell_width,
            cell_height=cell_height,
            slide_width=self.slide_width,
            slide_height=self.slide_height,
            rename={},
        )

    def add_summary_table_to_pptx(
        self,
        title,
        cell_width,
        items,
        cell_height=20,
        group_by_table=None,
        rename=None,
        pin_kind=None,
        sort=None,
        merge=False,
    ):
        """add summary table to pptx

        add slide, add table to pptx.

            Args:
                title (str): slide title
                cell_width (list): cell width
                items (list): items for table
                cell_height (int): cell height
                group_by_table (str): group name of table in case separate table by group
                rename (dict): specify header original and after name in case rename
                pin_kind (str): pin kind
                sort (bool): if sort data by one of data frame column, specify df column name.
                merge (bool): marge flag

            Returns:
                None

        """

        # this code works only for overshoot/undershoot
        if merge:
            new_data = []
            for i in range(0, len(self.data_vns), 2):
                new_dic = dict({**self.data_vns[i], **self.data_vns[i + 1]})
                new_data.append(new_dic)

            df_vns = pd.DataFrame(new_data)
            self.data_df = pd.merge(self.data_df, df_vns, on="Condition", how="outer")

        self.add_slide_to_pptx(title=title, layout=4)

        # if needs to separate result per pin kind
        if pin_kind:
            data_list_to_table_df = (
                self.data_df[self.data_df["Pin_kind"] == pin_kind].copy().reset_index()
            )
        else:
            data_list_to_table_df = self.data_df.reset_index()

        if sort:
            data_list_to_table_df = data_list_to_table_df.sort_values(sort)

        # data_list_to_table_df = self.data_df.reset_index()
        self.add_table(
            df=data_list_to_table_df,
            items=items,
            cell_width=cell_width,
            cell_height=cell_height,
            slide_width=self.slide_width,
            slide_height=self.slide_height,
            rename=rename,
        )

        if group_by_table:
            for name, group in data_list_to_table_df.groupby(group_by_table):
                self.add_slide_to_pptx(title=str(name), layout=4)

                data_list_to_table_df = group.reset_index()
                self.add_table(
                    df=data_list_to_table_df,
                    items=items,
                    cell_width=cell_width,
                    cell_height=cell_height,
                    slide_width=self.slide_width,
                    slide_height=self.slide_height,
                    rename=rename,
                )

    def add_pictures_to_pptx(self, *file_list, resize=False, picture_width=500):
        """add pictures to pptx

        Args:
            *file_list (list): picture list
            resize (bool): resize flag
            picture_width (int): picture width

        Returns:
            None

        """

        picture_per_slide = len(file_list)

        if self.pptx_lib == "win32com":
            picture_width = picture_width

        elif self.pptx_lib == "python-pptx":
            picture_width = Pt(picture_width)

        top = self.slide_height * 0.3
        left_1 = self.slide_width / 4 - picture_width / 2
        left_2 = self.slide_width * 3 / 4 - picture_width / 2
        text_box_height = 40

        if picture_per_slide == 1:
            for file in file_list[0]:
                # TODO change title
                title = os.path.splitext(os.path.basename(file))[0]
                for key, value in RENAME_CONDITIONS.items():
                    title = title.replace(key, value)

                self.add_slide_to_pptx(title=title, layout=11)
                self.add_picture_to_pptx(
                    file_path=file,
                    resize=resize,
                    picture_width=picture_width,
                    count_picture=False,
                )

        elif picture_per_slide == 2:
            for (file1, file2) in zip(file_list[0], file_list[1]):
                title = (
                    file1.replace("\\", "xyz")
                    .replace(".png", "")
                    .replace("8GPE_Frequency", "")  # TODO need to fix title
                )
                title = re.sub(".*xyz", "", title)
                self.add_slide_to_pptx(title=title, layout=11)

                title1 = file1.replace("\\", "xyz").replace(".png", "")
                title1 = re.sub(".*xyz", "", title1)

                # 1st picture
                self.add_picture_to_pptx(
                    file_path=file1,
                    count_picture=False,
                    picture_width=picture_width,
                    resize=resize,
                    reposition=True,
                    top=top,
                    left=left_1,
                )
                # 1st text box
                self.add_textbox(
                    title=title1,
                    slide_num=self.slide_count,
                    left=left_1,
                    top=top,
                    width=picture_width,
                    height=text_box_height,
                )

                title2 = file2.replace("\\", "xyz").replace(".png", "")
                title2 = re.sub(".*xyz", "", title2)

                # 2nd picture
                self.add_picture_to_pptx(
                    file_path=file2,
                    count_picture=False,
                    picture_width=picture_width,
                    resize=resize,
                    reposition=True,
                    top=top,
                    left=left_2,
                )
                # 2nd text box
                self.add_textbox(
                    title=title2,
                    slide_num=self.slide_count,
                    left=left_2,
                    top=top,
                    width=picture_width,
                    height=text_box_height,
                )

    def add_picture_to_pptx(
        self,
        file_path,
        count_picture=True,
        picture_width=400,
        left=0,
        resize=False,
        reposition=False,
        top=0,
    ):
        """add picture to pptx

        Args:
            file_path (str): file name full path
            count_picture (bool): picture counter
            picture_width (float): picture width
            resize (bool): set True if resize picture
            reposition (bool): picture reposition flag
            top (float): picture top position
            left (float): picture left position

        Returns:
            None

        """
        global picture_counter

        picture = None
        if self.pptx_lib == "win32com":
            picture = self.active_presentation.Slides(
                self.slide_count
            ).Shapes.AddPicture(
                FileName=file_path, LinkToFile=-1, SaveWithDocument=-1, Left=0, Top=0,
            )

        elif self.pptx_lib == "python-pptx":
            # im = Image.open(file_path)
            # im_width, im_height = im.size

            if resize:
                picture = self.slide.shapes.add_picture(
                    image_file=file_path, left=0, top=0, width=picture_width
                )
            else:
                picture = self.slide.shapes.add_picture(
                    image_file=file_path, left=0, top=0
                )

        if self.pptx_lib == "win32com":
            if resize:
                picture.Width = picture_width

            if reposition:
                picture.Top = top
                picture.Left = left

            else:
                picture.Top = self.slide_height / 2 - picture.Height / 2
                picture.Left = self.slide_width / 2 - picture.Width / 2

        elif self.pptx_lib == "python-pptx":

            if reposition:
                picture.top = int(top)
                picture.left = int(left)

            else:
                picture.top = int(self.slide_height / 2 - picture.height / 2)
                picture.left = int(self.slide_width / 2 - picture.width / 2)

        if count_picture:
            picture_counter += 1

    def add_textbox(self, title, slide_num, top, left, width, height, font_size=14):
        if self.pptx_lib == "win32com":
            text_box = self.active_presentation.Slides(slide_num).Shapes.AddTextbox(
                1, Top=top, Left=left, Width=width, Height=height,
            )
            text_box.TextFrame.TextRange.Text = title
            text_box.TextFrame.TextRange.ParagraphFormat.Alignment = 2  # center
            text_box.TextFrame.VerticalAnchor = 4  # bottom
            text_box.TextFrame.TextRange.Font.Size = font_size
            text_box.Top = text_box.Top - text_box.Height
            text_box.Height = height
            text_box.Top = top - height

        elif self.pptx_lib == "python-pptx":
            text_box = self.slide.shapes.add_textbox(
                left=left, top=top - Pt(height), width=width, height=Pt(height)
            )
            text_frame = text_box.text_frame
            text_frame.paragraphs[0].font.size = Pt(font_size)
            pg = text_frame.paragraphs[0]
            pg.text = title.replace(".png", "").replace(".PNG", "")
            pg.alignment = PP_ALIGN.CENTER
            text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM

    def add_slide_to_pptx(self, title, layout, font_size=20):
        """add slide to pptx

        Args:
            title (str): slide title
            layout (int): slide layout
            font_size (int): font size

        Returns:
            None

        """
        self.update_slide_count()

        if self.pptx_lib == "win32com":
            self.slide = self.active_presentation.Slides.Add(
                Index=self.slide_count + 1, Layout=layout
            )
            self.slide.Select()
            self.slide.Shapes(1).TextFrame.TextRange.Text = title
            self.slide.Shapes(1).TextFrame.TextRange.Font.Size = font_size
            self.slide_count += 1

        elif self.pptx_lib == "python-pptx":
            # TODO check slide layout for table
            self.slide = self.active_presentation.slides.add_slide(
                self.active_presentation.slide_layouts[6]
            )
            self.slide.shapes[0].text = title
            self.slide.shapes[0].text_frame.paragraphs[0].font.size = Pt(font_size)

    def add_table(
        self, df, items, cell_width, cell_height, slide_width, slide_height, rename,
    ):
        """add table to slide.

        Args:
            df (DataFrame): data to table. data type is pandas dataframe
            items (list): items to add table
            cell_width (list): cell width
            cell_height (int): cell height
            slide_width (int): slide width
            slide_height (int): slide height
            rename (dict): specify before/after name as dict like {"Frequency":"Freq(GHz)"} if rename table header.

        Returns:
            None

        """
        df = df.loc[:, items]
        print(df)
        df = df.dropna(how="all", axis=1)  # drop all nan column
        data_list_to_table = df.values.tolist()
        data_list_to_table.insert(0, df.columns.tolist())

        table_rows = len(data_list_to_table)
        table_columns = len(data_list_to_table[0])
        table = None
        tr = None
        table_shape = None
        if self.pptx_lib == "win32com":
            table = self.slide.Shapes.AddTable(table_rows, table_columns).Table

        elif self.pptx_lib == "python-pptx":
            table_shape = self.slide.shapes.add_table(
                table_rows,
                table_columns,
                left=0,
                top=0,
                width=self.slide_width,
                height=self.slide_height,
            )

            table = table_shape.table

        for i in range(table_rows):
            for j in range(table_columns):
                if self.pptx_lib == "win32com":
                    tr = table.Cell(i + 1, j + 1).Shape.TextFrame.TextRange

                elif self.pptx_lib == "python-pptx":
                    tr = table.cell(i, j)

                try:
                    if (
                        data_list_to_table[i][j] == 9.91e37
                        or data_list_to_table[i][j] == 9.91e40
                    ):
                        data_list_to_table[i][j] = "acquisition failure"

                    elif str(data_list_to_table[i][j]) == "nan":
                        data_list_to_table[i][j] = "-"

                    if self.pptx_lib == "win32com":
                        if "v-ns" in data_list_to_table[0][j]:
                            tr.Text = f"{data_list_to_table[i][j]:.6f}"
                        else:
                            tr.Text = f"{data_list_to_table[i][j]:.1f}"

                    elif self.pptx_lib == "python-pptx":
                        if "v-ns" in data_list_to_table[0][j]:
                            tr.text = f"{data_list_to_table[i][j]:.6f}"
                        else:
                            tr.text = f"{data_list_to_table[i][j]:.1f}"

                except ValueError:
                    if rename:
                        if data_list_to_table[i][j] in rename:
                            if self.pptx_lib == "win32com":
                                tr.Text = rename[data_list_to_table[i][j]]

                            elif self.pptx_lib == "python-pptx":
                                tr.text = rename[data_list_to_table[i][j]]

                        else:
                            if self.pptx_lib == "win32com":
                                tr.Text = data_list_to_table[i][j]

                                if data_list_to_table[0][j] == "Vi":
                                    for key, value in RENAME_CONDITIONS.items():
                                        tr.Text = tr.Text.replace(key, value)

                            elif self.pptx_lib == "python-pptx":
                                tr.text = data_list_to_table[i][j]

                                if data_list_to_table[0][j] == "Vi":
                                    for key, value in RENAME_CONDITIONS.items():
                                        tr.text = tr.Text.replace(key, value)

                    else:
                        if self.pptx_lib == "win32com":
                            tr.Text = data_list_to_table[i][j]

                            if data_list_to_table[0][j] == "Vi":
                                for key, value in RENAME_CONDITIONS.items():
                                    tr.Text = tr.Text.replace(key, value)

                        elif self.pptx_lib == "python-pptx":
                            tr.text = data_list_to_table[i][j]

                            if data_list_to_table[0][j] == "Vi":
                                for key, value in RENAME_CONDITIONS.items():
                                    tr.text = tr.Text.replace(key, value)

                if self.pptx_lib == "win32com":
                    tr.Font.Size = 10
                    tr.ParagraphFormat.Alignment = 2  # centering

                elif self.pptx_lib == "python-pptx":
                    # TODO this doesn't work. why?
                    # pf = tr.text_frame.paragraphs
                    # for i in range(len(pf)):
                    #     pf[i].font.size = Pt(10)

                    # This works
                    tr.text_frame.paragraphs[0].font.size = Pt(10)
                    tr.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    if len(tr.text_frame.paragraphs) == 2:
                        tr.text_frame.paragraphs[1].font.size = Pt(10)
                        tr.text_frame.paragraphs[1].alignment = PP_ALIGN.CENTER

        if self.pptx_lib == "win32com":
            for i in range(1, table_columns + 1):
                table.Columns(i).Width = cell_width[i - 1]

            for i in range(1, table_rows + 1):
                table.Rows(i).Height = cell_height

            # adjust table position
            shape = self.slide.Shapes(2)
            shape.Left = slide_width / 2 - shape.width / 2
            shape.Top = slide_height / 6

        elif self.pptx_lib == "python-pptx":
            for i in range(table_columns):
                table.columns[i].width = Pt(cell_width[i])

            for i in range(table_rows):
                table.rows[i].height = Pt(cell_height)

            # adjust table position
            table_shape.left = int(slide_width / 2 - table_shape.width / 2)
            table_shape.top = int(slide_height / 6)

    def save_pptx(self, file_name, folder_name):
        """save pptx file

        Args:
            file_name (str): file name
            folder_name (str): folder name

        Returns:
            None

        """
        file_path = folder_name + str(date_now) + "_" + file_name
        if self.pptx_lib == "win32com":
            self.active_presentation.SaveAs(FileName=file_path)

        elif self.pptx_lib == "python-pptx":
            self.active_presentation.save(file=file_path)

    def adjust_unit(self):
        """adjust unit of data

        Args:

        Returns:
            None

        """
        if "Eheight" in self.data_df.columns:
            self.data_df["Eheight"] = self.data_df["Eheight"].apply(mul3)

        if "Ewidth" in self.data_df.columns:
            self.data_df["Ewidth"] = self.data_df["Ewidth"].apply(mul12)

        if "Risetime" in self.data_df.columns:
            self.data_df["Risetime"] = self.data_df["Risetime"].apply(mul12)

        if "Falltime" in self.data_df.columns:
            self.data_df["Falltime"] = self.data_df["Falltime"].apply(mul12)

        if "Frequency" in self.data_df.columns:
            self.data_df["Frequency"] = self.data_df["Frequency"].apply(mul_m9)

        if "Vamplitude" in self.data_df.columns:
            self.data_df["Vamplitude"] = self.data_df["Vamplitude"].apply(mul3)

        if "Vpp" in self.data_df.columns:
            self.data_df["Vpp"] = self.data_df["Vpp"].apply(mul3)

        if "Vmaximum" in self.data_df.columns:
            self.data_df["Vmaximum"] = self.data_df["Vmaximum"].apply(mul3)

        if "Vminimum" in self.data_df.columns:
            self.data_df["Vminimum"] = self.data_df["Vminimum"].apply(mul3)

        if "Vtop" in self.data_df.columns:
            self.data_df["Vtop"] = self.data_df["Vtop"].apply(mul3)

        if "Top" in self.data_df.columns:
            self.data_df["Top"] = self.data_df["Top"].apply(mul3)

        if "Vbase" in self.data_df.columns:
            self.data_df["Vbase"] = self.data_df["Vbase"].apply(mul3)

        if "Base" in self.data_df.columns:
            self.data_df["Base"] = self.data_df["Base"].apply(mul3)

        if "Pwidth" in self.data_df.columns:
            self.data_df["Pwidth"] = self.data_df["Pwidth"].apply(mul12)

        if "Pp" in self.data_df.columns:
            self.data_df["Pp"] = self.data_df["Pp"].apply(mul12)


if __name__ == "__main__":
    start = time.time()
    DATA_START_COLUMNS = 10
    FOLDER_PATH = os.getcwd() + "/202106172016_debug/"
    PPTX_FILE_NAME = "8GPE_TEST.pptx"
    OSC_PICTURE_LIST_CROSSTALK = glob(FOLDER_PATH + "/*crosstalk/*.png")

    PE = "8GPE_"
    # PIN_KINDS = ["IO", "WCK", "CK", "CA", "CS"]
    PIN_KINDS = ["IO"]
    pin_kind_for_pptx = "IO"
    PPTX_LIB = "win32com"
    # PPTX_LIB = "python-pptx"

    if PPTX_LIB == "win32com":
        pptx = win32com.client.Dispatch("PowerPoint.Application")
        pptx.Visible = True
        active_presentation_object = pptx.Presentations.Open(
            os.getcwd() + "/advtemplate_mini.pptx"
        )

    else:
        active_presentation_object = Presentation(
            os.getcwd() + "/advtemplate_mini.pptx"
        )

    OVERSHOOT_FILE_LIST = glob(
        FOLDER_PATH + pin_kind_for_pptx.lower() + "/*_overshoot/*.txt"
    )
    OSC_PICTURE_LIST_OVERSHOOT = glob(
        FOLDER_PATH + pin_kind_for_pptx.lower() + "/*overshoot/*.png"
    )
    OSC_PICTURE_LIST_VIH_DC = glob(
        FOLDER_PATH + pin_kind_for_pptx.lower() + "/*vih_dc/*.png"
    )
    OSC_PICTURE_LIST_VIL_DC = glob(
        FOLDER_PATH + pin_kind_for_pptx.lower() + "/*vil_dc/*.png"
    )
    OSC_PICTURE_LIST_VIHL_AC = glob(
        FOLDER_PATH + pin_kind_for_pptx.lower() + "/*vihl_ac/*.png"
    )
    # Overshoot/Undershoot
    # wave_data_overshoot = WaveData(
    #     active_presentation_object=active_presentation_object,
    #     file_name=pin_kind_for_pptx.lower() + "_" + OVERSHOOT_FILE_NAME,
    #     folder_path=FOLDER_PATH + pin_kind_for_pptx.lower() + "/",
    #     group_by=DATA_GROUP,
    #     index=DATA_INDEX,
    #     pptx_lib=PPTX_LIB,
    # )
    # for i in range(len(OVERSHOOT_FILE_LIST)):
    #     if i < 4:
    #         y_ticks = OVERSHOOT_Y_TICKS_1V_0V
    #     else:
    #         y_ticks = OVERSHOOT_Y_TICKS_0r5V_0V

    #     if i % 2 == 0:
    #         item_name = "Overshoot"
    #     else:
    #         item_name = "Undershoot"

    #     wave_data_overshoot.make_overshoot_graph(
    #         file=OVERSHOOT_FILE_LIST[i],
    #         y_ticks=y_ticks,
    #         item_name=item_name,
    #         additional_information=True,
    #         info="Target: under 300mV, under 0.1 V-ns",
    #     )
    # wave_data_overshoot.add_summary_table_to_pptx(
    #     title="Overshoot/Undershoot",
    #     cell_width=[
    #         CELL_WIDTH_BASE * 0.9,  # pin
    #         CELL_WIDTH_BASE * 2,  # vi
    #         CELL_WIDTH_BASE * 0.9,  # rate
    #         CELL_WIDTH_BASE * 1.2,  # overshoot
    #         CELL_WIDTH_BASE * 1.54,  # v-ns-overshoot
    #         CELL_WIDTH_BASE * 0.9,  # Vmaximum
    #         CELL_WIDTH_BASE * 0.9,  # vtop
    #         CELL_WIDTH_BASE * 1.2,  # undershoot
    #         CELL_WIDTH_BASE * 1.56,  # v-ns-undershoot
    #         CELL_WIDTH_BASE * 0.9,  # Vminimum
    #         CELL_WIDTH_BASE * 0.95,  # base
    #     ],
    #     items=[
    #         "Pin",
    #         "Vi",
    #         "Rate",
    #         "Overshoot_mV",
    #         "v-ns-Overshoot",
    #         "Vmaximum",
    #         "Vtop",
    #         "Undershoot_mV",
    #         "v-ns-Undershoot",
    #         "Vminimum",
    #         "Vbase",
    #     ],
    #     group_by_table="Vi",
    #     rename={
    #         "Overshoot_mV": "Overshoot(mV)",
    #         "Vmaximum": "Vmax(mV)",
    #         "Vtop": "Vtop(mv)",
    #         "Undershoot_mV": "Undershoot(mV)",
    #         "Vminimum": "Vmin(mV)",
    #         "Vbase": "Vbase(mV)",
    #         "v-ns-Overshoot": "v-ns-Overshoot\n(V-ns)",
    #         "v-ns-Undershoot": "v-ns-Undershoot\n(V-ns)",
    #         "nan": "-",
    #     },
    #     merge=True,
    #     # sort="Order"
    #     # pin_kind_for_pptx="IO"
    # )
    # wave_data_overshoot.add_pictures_to_pptx(
    #     OSC_PICTURE_LIST_OVERSHOOT,
    #     resize=True,
    # )
    wave_data_overview = WaveData(
        active_presentation=active_presentation_object,
        file_name=pin_kind_for_pptx.lower() + "_" + OVERVIEW_FILE_NAME,
        folder_path=FOLDER_PATH + pin_kind_for_pptx.lower() + "/",
        group_by=DATA_GROUP,
        index=DATA_INDEX,
        pptx_lib=PPTX_LIB,
    )
    wave_data_overview.make_vix_graph(
        item_name="vix",
        positive_pin_file=FOLDER_PATH
        + "P1857A1_overview_Vih1r000V_Vil0r000V_Vt0r500V_Rate0r286ns_Duty0r500.txt",
        negative_pin_file=FOLDER_PATH
        + "P1858A1_overview_Vih1r000V_Vil0r000V_Vt0r500V_Rate0r286ns_Duty0r500.txt",
    )
    for pin_kind_for_pptx in PIN_KINDS:
        wave_data_overview.make_graph(
            df_columns_list=["Frequency"],
            file_name=PE + "Frequency",
            digit_format="%.2f",
            legends={"Frequency": "Freq(GHz)"},
            ax_h_lines=[4.0],
            y_ticks=FREQ_YTICKS,
            y_label="GHz",
            pin_kind=pin_kind_for_pptx,
            ax_h_lines_per_condition={"IO_Vih1r3V_Vil-0r50V_Vt-0r50V": [2.0, 3.0]},
            y_ticks_per_condition={"CS_Vih1r0V_Vil0r0V_Vt0r5V": [0, 2, 0.2]},
            additional_information=True,
            info="Spec: less than 60ps (@1.0Vp-p/20% to 80%)",
        )
        wave_data_overview.add_summary_table_to_pptx(
            title="overview_" + pin_kind_for_pptx,
            cell_width=[
                CELL_WIDTH_BASE_PIN * 1.1,
                CELL_WIDTH_BASE_VI * 2.0,
                CELL_WIDTH_BASE_RATE * 1.1,
                CELL_WIDTH_BASE * 1.1,
            ],
            items=["Pin", "Vi", "Rate", "Frequency"],
            rename={"Vih1r0V_Vil0r0V_Vt0r5V": "Vih=1.0V, Vil=0.0V"},
            # rename={"Vih1r0V_": "Vih=1.0V,"},
            pin_kind=pin_kind_for_pptx,
        )
    # wave_data_vih = WaveData(
    #     active_presentation_object=active_presentation_object,
    #     file_name=pin_kind_for_pptx.lower() + "_" + VIH_DC_FILE_NAME,
    #     folder_path=FOLDER_PATH + pin_kind_for_pptx.lower() + "/",
    #     group_by=DATA_GROUP,
    #     index=DATA_INDEX,
    #     pptx_lib=PPTX_LIB,
    # )
    # wave_data_vih.make_graph(
    #     df_columns_list=["Vmaximum", "Vminimum"],
    #     file_name=PE + "Vih_Dc",
    #     digit_format="%.1f",
    #     legends=["Vmax(mV)", "Vmin(mV)"],
    #     # ax_h_lines=[4.0],
    #     y_ticks=[100, 1100, 100],
    #     y_label="mV",
    #     pin_kind_for_pptx=pin_kind_for_pptx,
    # )
    # wave_data_vih.add_summary_table_to_pptx(
    #     title="Vih Dc summary",
    #     cell_width=[
    #         CELL_WIDTH_BASE_PIN,
    #         CELL_WIDTH_BASE_VI,
    #         CELL_WIDTH_BASE_RATE,
    #         CELL_WIDTH_BASE * 2,
    #         CELL_WIDTH_BASE * 2,
    #     ],
    #     items=["Pin", "Vi", "Rate", "Vmaximum", "Vminimum"],
    #     rename={"Vmaximum": "Vmax(mV)", "Vminimum": "Vmin(mV)"},
    #     # rename={"Vih1r0V_": "Vih=1.0V,"},
    #     pin_kind_for_pptx=pin_kind_for_pptx,
    # )
    # wave_data_vih.add_pictures_to_pptx(
    #     OSC_PICTURE_LIST_VIH_DC, resize=True,
    # )
    # wave_data_vil = WaveData(
    #     active_presentation_object=active_presentation_object,
    #     file_name=pin_kind_for_pptx.lower() + "_" + VIL_DC_FILE_NAME,
    #     folder_path=FOLDER_PATH + pin_kind_for_pptx.lower() + "/",
    #     group_by=DATA_GROUP,
    #     index=DATA_INDEX,
    #     pptx_lib=PPTX_LIB,
    # )
    # wave_data_vil.make_graph(
    #     df_columns_list=["Vmaximum", "Vminimum"],
    #     file_name=PE + "vil_dc",
    #     digit_format="%.1f",
    #     legends=["Vmax(mV)", "Vmin(mV)"],
    #     # ax_h_lines=[4.0],
    #     y_ticks=[-300, 700, 100],
    #     y_label="mV",
    #     pin_kind_for_pptx=pin_kind_for_pptx,
    # )
    # wave_data_vil.add_summary_table_to_pptx(
    #     title="VIL_DC summary",
    #     cell_width=[
    #         CELL_WIDTH_BASE_PIN,
    #         CELL_WIDTH_BASE_VI,
    #         CELL_WIDTH_BASE_RATE,
    #         CELL_WIDTH_BASE * 2,
    #         CELL_WIDTH_BASE * 2,
    #     ],
    #     items=["Pin", "Vi", "Rate", "Vmaximum", "Vminimum"],
    #     rename={"Vmaximum": "Vmax(mV)", "Vminimum": "Vmin(mV)"},
    #     # rename={"Vih1r0V_": "Vih=1.0V,"},
    #     pin_kind_for_pptx=pin_kind_for_pptx,
    # )
    # wave_data_vil.add_pictures_to_pptx(
    #     OSC_PICTURE_LIST_VIL_DC, resize=True,
    # )
    # wave_data_vihl_ac = WaveData(
    #     active_presentation_object=active_presentation_object,
    #     file_name=pin_kind_for_pptx.lower() + "_" + VIHL_AC_FILE_NAME,
    #     folder_path=FOLDER_PATH + pin_kind_for_pptx.lower() + "/",
    #     group_by=DATA_GROUP,
    #     index=DATA_INDEX,
    #     pptx_lib=PPTX_LIB,
    # )
    # wave_data_vihl_ac.make_graph(
    #     df_columns_list=["Vmaximum", "Vminimum"],
    #     file_name=PE + "vihl_ac",
    #     digit_format="%.1f",
    #     legends=["Vmax(mV)", "Vmin(mV)"],
    #     # ax_h_lines=[4.0],
    #     y_ticks=[-300, 1100, 100],
    #     y_label="mV",
    #     pin_kind_for_pptx=pin_kind_for_pptx,
    # )
    # wave_data_vihl_ac.add_summary_table_to_pptx(
    #     title="VIHL_AC summary",
    #     cell_width=[
    #         CELL_WIDTH_BASE_PIN,
    #         CELL_WIDTH_BASE_VI,
    #         CELL_WIDTH_BASE_RATE,
    #         CELL_WIDTH_BASE * 2,
    #         CELL_WIDTH_BASE * 2,
    #     ],
    #     items=["Pin", "Vi", "Rate", "Vmaximum", "Vminimum"],
    #     rename={"Vmaximum": "Vmax(mV)", "Vminimum": "Vmin(mV)"},
    #     # rename={"Vih1r0V_": "Vih=1.0V,"},
    #     pin_kind_for_pptx=pin_kind_for_pptx,
    # )
    # wave_data_vihl_ac.add_pictures_to_pptx(
    #     OSC_PICTURE_LIST_VIHL_AC, resize=True,
    # )
    wave_data_overview.save_pptx(file_name=PPTX_FILE_NAME, folder_name=FOLDER_PATH)
    elapsed_time = time.time() - start
    print(f"exec time:{elapsed_time:.1f}[sec]")
