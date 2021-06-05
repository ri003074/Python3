import sys

abc = ["abc", "def", "ghi"]

dic = {"abc": "AAA", "def": "DDD"}


for i in range(len(abc)):
    for key, value in dic.items():
        if abc[i] == key:
            abc[i] = value

print(abc)


class Outer:
    def __init__(self):
        print("create Outer Class")
        print(self.Inner)
        # self.Innerのメモリサイズを確認
        print(sys.getsizeof(self.Inner))
        self.inner = self.Inner()
        print(self.inner)

    class Inner:
        def __init__(self):
            print("create Inner Class")

        def abc(self):
            print("inner abc")


outer = Outer()

outer.inner.abc()

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN


prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = Cm(3)
top = Cm(2.5)
width = Cm(15)
height = Cm(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "Hello"
txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Just an example"
font = run.font


prs.save("test.pptx")

