import os
import sys
import time

import win32com

from data_summarize import WaveData
from glob import glob
from pptx import Presentation

if __name__ == "__main__":
    start = time.time()
    CELL_WIDTH_BASE = 72
    DATA_START_COLUMNS = 10
    FOLDER_PATH = os.getcwd() + "/20210602_debug_8gpe/"
    PPTX_FILE_NAME = "8GPE_TEST.pptx"
    OVERVIEW_FILE_NAME = "result_overview.csv"
    EYE_FILE_NAME = "result_eye.csv"
    HISTOGRAM_FILE_NAME = "result_histogram.csv"
    OSC_PICTURE_LIST_OVERVIEW = glob(FOLDER_PATH + "/*overview/*.png")
    OSC_PICTURE_LIST_EYE = glob(FOLDER_PATH + "/*eye/*.png")
    OSC_PICTURE_LIST_HISTOGRAM = glob(FOLDER_PATH + "/*histogram/*.png")
    DATA_GROUP = "Pkind_Vi"
    DATA_INDEX = "Pin_Rate"
    FREQ_YTICKS = [1.0, 5.0, 0.5]
    # DUTY_YTICKS = [40.0, 60.0, 2.5]
    DUTY_YTICKS = [41.0, 59.0, 3.0]
    TRTF_YTICKS = [30.0, 70.0, 5]
    EHEIGHT_YTICS = [300, 400, 20]
    EWIDTH_YTICKS = [60, 120, 10]
    PP_YTICKS = [00, 50, 10]

    PE = "8GPE_"
    PKINDS = ["IO", "WCK", "CK", "CA", "CS"]
    pkind = "IO"
    # PPTX_LIB = "win32com"
    PPTX_LIB = "python-pptx"

    if PPTX_LIB == "win32com":
        pptx = win32com.client.Dispatch("PowerPoint.Application")
        pptx.Visible = True
        active_presentation = pptx.Presentations.Open(
            os.getcwd() + "/advtemplate_mini.pptx"
        )

    elif PPTX_LIB == "python-pptx":
        active_presentation = Presentation(os.getcwd() + "/advtemplate_mini.pptx")

    wave_data_overview = WaveData(
        active_presentation=active_presentation,
        file_name=OVERVIEW_FILE_NAME,
        folder_path=FOLDER_PATH,
        groupby=DATA_GROUP,
        index=DATA_INDEX,
        pptx_lib=PPTX_LIB,
    )
    # wave_data_overview.make_overshoot_graph(
    #     file=FOLDER_PATH
    #     + "20210602_101849_P111A1_overview/P111A1_overview_Vih0r500V_Vil0r000V_Vt0r000V_Rate0r286ns_Duty0r500.txt"
    # )
    # wave_data_overview.make_vix_graph(
    #     posi_pin_file="./sample_log/P1859A2_RZX_Vih0r916V_Vil0r000V_Rate0r362ns_Speed2r759GHz_TopBase_Meas.txt",
    #     nega_pin_file="./sample_log/P1860A2_RZX_Vih0r916V_Vil0r000V_Rate0r362ns_Speed2r759GHz_TopBase_Meas.txt",
    #     description=True,
    #     item_name=PE + "Vix",
    #     reference_level=0.2,
    #     ylabel="mV",
    # )
    # wave_data_overview.make_vix_graph(
    #     posi_pin_file="./sample_log/P1859A2_RZX_Vih0r916V_Vil0r000V_Rate0r362ns_Speed2r759GHz_TopBase_Meas.txt",
    #     nega_pin_file="./sample_log/P1860A2_RZX_Vih0r916V_Vil0r000V_Rate0r362ns_Speed2r759GHz_TopBase_Meas.txt",
    #     description=False,
    #     item_name=PE + "Vix",
    #     reference_level=0.229,
    #     ylabel="mV",
    # )
    # wave_data_overview.add_vix_table_to_pptx(
    #     title="Vix",
    #     items=[
    #         "Positive Pin",
    #         "Negative Pin",
    #         "Vi",
    #         "rate",
    #         "Vix_WCK_FR/|Min(f(t))| (%)",
    #         "Vix_WCK_Rf/Max(f(t)) (%)",
    #     ],
    #     cell_width=[
    #         CELL_WIDTH_BASE * 1.1,
    #         CELL_WIDTH_BASE * 1.1,
    #         CELL_WIDTH_BASE * 2.0,
    #         CELL_WIDTH_BASE * 1.1,
    #         CELL_WIDTH_BASE * 2.0,
    #         CELL_WIDTH_BASE * 2.0,
    #     ],
    #     cell_height=20,
    # )
    # for pkind in PKINDS:
    wave_data_overview.make_graph(
        df_columns_list=["Frequency"],
        file_name=PE + "Frequency",
        format="%.2f",
        legends=["Freq(GHz)"],
        legend_loc="lower right",
        pkind=pkind,
        yticks=FREQ_YTICKS,
        ylabel="GHz",
    )
    wave_data_overview.make_graph(
        axhline=[47, 53],  # reference line
        df_columns_list=["Dutycycle"],
        file_name=PE + "Duty",
        legends=["Duty(%)"],
        pkind=pkind,
        yticks=DUTY_YTICKS,
        ylabel="%",
    )
    wave_data_overview.make_graph(
        axhline=[60],  # spec line
        df_columns_list=["Risetime", "Falltime"],
        file_name=PE + "Risetime_Falltime",
        legends=["Tr(ps)", "Tf(ps)"],
        pkind=pkind,
        spec=True,
        yticks=TRTF_YTICKS,
        ylabel="ps",
    )
    wave_data_overview.add_summary_table_to_pptx(
        title="overview",
        cell_width=[
            CELL_WIDTH_BASE * 1.1,
            CELL_WIDTH_BASE * 2,
            CELL_WIDTH_BASE * 1.2,
            CELL_WIDTH_BASE * 1.2,
            CELL_WIDTH_BASE * 1.2,
            CELL_WIDTH_BASE * 1.2,
            CELL_WIDTH_BASE * 1.2,
            CELL_WIDTH_BASE * 1.2,
            CELL_WIDTH_BASE * 1.2,
            CELL_WIDTH_BASE * 1.2,
            CELL_WIDTH_BASE * 1.2,
            CELL_WIDTH_BASE * 1.2,
        ],
        items=[
            "Pin",
            "Vi",
            "Rate",
            "Frequency",
            "Dutycycle",
            "Risetime",
            "Falltime",
            "Overshoot",
            "Preshoot",
            "Pwidth",
        ],
        groupby_table="Vi",
        rename={
            "Risetime": "Tr(ps)",
            "Frequency": "Freq(GHz)",
            "Dutycycle": "Duty(%)",
            "Falltime": "Tf(ps)",
            "Overshoot": "Overshoot(%)",
            "Preshoot": "Preshoot(%)",
            "Pwidth": "Pwidth(ps)",
            "nan": "-",
        },
        # sort="Order"
        pkind=pkind,
    )
    wave_data_overview.save_pptx(file_name=PPTX_FILE_NAME, folder_name=FOLDER_PATH)
    elapsed_time = time.time() - start
    print(f"elapsed_time:{elapsed_time:.1f}[sec]")
    sys.exit()
    wave_data_eye = WaveData(
        active_presentation=active_presentation,
        file_name=EYE_FILE_NAME,
        folder_path=FOLDER_PATH,
        groupby=DATA_GROUP,
        index=DATA_INDEX,
        pptx_lib=PPTX_LIB,
    )
    wave_data_eye.make_graph(
        df_columns_list=["Eheight"],
        file_name=PE + "Eheight",
        legends=["Eye Height(mV)"],
        ylabel="mV",
        yticks=EHEIGHT_YTICS,
    )
    wave_data_eye.add_summary_table_to_pptx(
        title="eye",
        cell_width=[
            # CELL_WIDTH_BASE * 5,
            CELL_WIDTH_BASE * 2,
            CELL_WIDTH_BASE * 2,
            CELL_WIDTH_BASE * 2,
            CELL_WIDTH_BASE * 2,
            CELL_WIDTH_BASE * 2,
            CELL_WIDTH_BASE * 2,
            CELL_WIDTH_BASE * 2,
        ],
        items=["Pin", "Vi", "Rate", "Eheight"],
        rename={"Eheight": "Eye Height(mV)"},
        groupby_table="Vi",
    )
    # wave_data_eye.make_excel_graphs(
    #     data_start_column=DATA_START_COLUMNS,
    #     chart_yaxis_titles=["ps", "mV"],
    #     chart_yaxis_scaling_mins=[300, 0],
    #     chart_yaxis_scaling_maxes=[400, 10],
    #     chart_yaxis_major_unit=[20, 2],
    # )
    wave_data_histogram = WaveData(
        active_presentation=active_presentation,
        file_name=HISTOGRAM_FILE_NAME,
        folder_path=FOLDER_PATH,
        groupby=DATA_GROUP,
        index=DATA_INDEX,
        pptx_lib=PPTX_LIB,
    )
    wave_data_histogram.make_graph(
        df_columns_list=["Pp"],
        file_name=PE + "Jitter",
        legends=["PP(ps)"],
        ylabel="ps",
        yticks=PP_YTICKS,
    )
    wave_data_histogram.add_pictures_to_pptx(file_list=OSC_PICTURE_LIST_OVERVIEW,)
    wave_data_histogram.add_pictures_to_pptx(file_list=OSC_PICTURE_LIST_EYE,)
    wave_data_histogram.add_pictures_to_pptx(file_list=OSC_PICTURE_LIST_HISTOGRAM,)
    wave_data_overview.save_pptx(file_name=PPTX_FILE_NAME, folder_name=FOLDER_PATH)
    elapsed_time = time.time() - start
    print(f"elapsed_time:{elapsed_time:.1f}[sec]")
